# # slide_renderer.py
# import os
# import uuid
# import pythoncom
# import win32com.client
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE


# def export_slide_to_png(ppt_path, slide_index):
#     pythoncom.CoInitialize()
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True

#     pres = powerpoint.Presentations.Open(ppt_path, WithWindow=True)
#     slide = pres.Slides[slide_index]

#     out_path = os.path.join(
#         os.path.dirname(ppt_path),
#         f"slide_{slide_index}_{uuid.uuid4().hex[:6]}.png"
#     )

#     slide.Export(out_path, "PNG", 1920, 1080)

#     pres.Close()
#     powerpoint.Quit()
#     pythoncom.CoUninitialize()

#     return out_path


# def _is_editable_text_shape(shape):
#     if not shape.has_text_frame:
#         return False

#     text = shape.text.strip()
#     if not text:
#         return False

#     if len(text) < 3:
#         return False

#     return True


# def extract_slide_structure(ppt_path, slide_index):
#     prs = Presentation(ppt_path)
#     slide = prs.slides[slide_index]

#     editable_shapes = []
#     idx = 0

#     for shape in slide.shapes:

#         # -----------------------------
#         # TEXT SHAPES
#         # -----------------------------
#         if _is_editable_text_shape(shape):

#             # ✅ SAFE placeholder detection
#             is_placeholder = False
#             try:
#                 _ = shape.placeholder_format
#                 is_placeholder = True
#             except Exception:
#                 is_placeholder = False

#             shape_entry = {
#                 "shape_id": f"shape_{idx}",
#                 "text": shape.text.strip(),
#                 "placeholder": is_placeholder,
#                 "type": "title" if is_placeholder else "body"
#             }

#             editable_shapes.append(shape_entry)
#             idx += 1

#         # -----------------------------
#         # GROUP SHAPES
#         # -----------------------------
#         elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
#             for shp in shape.shapes:
#                 if _is_editable_text_shape(shp):
#                     editable_shapes.append({
#                         "shape_id": f"shape_{idx}",
#                         "text": shp.text.strip(),
#                         "placeholder": False,
#                         "type": "body"
#                     })
#                     idx += 1

#     png_path = export_slide_to_png(ppt_path, slide_index)

#     return {
#         "slide_index": slide_index,
#         "ppt_path": ppt_path,
#         "png_path": png_path,
#         "editable_shapes": editable_shapes
#     }



# import os
# import uuid
# import platform
# import subprocess
# import shutil
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# # Detect OS
# IS_WINDOWS = platform.system() == "Windows"

# # =========================================================
# # PUBLIC API (USED BY FRONTEND)
# # =========================================================
# def extract_slide_structure(ppt_path, slide_index):
#    """
#    Extract editable text shapes from a slide and generate preview image.
#    Works on:
#    - Windows: PowerPoint COM
#    - Linux/Docker: LibreOffice
#    """
#    prs = Presentation(ppt_path)
#    slide = prs.slides[slide_index]
#    editable_shapes = []
#    idx = 0
#    for shape in slide.shapes:
#        # -----------------------------
#        # TEXT SHAPES
#        # -----------------------------
#        if _is_editable_text_shape(shape):
#            is_placeholder = False
#            try:
#                _ = shape.placeholder_format
#                is_placeholder = True
#            except Exception:
#                pass
#            editable_shapes.append({
#                "shape_id": f"shape_{idx}",
#                "text": shape.text.strip(),
#                "placeholder": is_placeholder,
#                "type": "title" if is_placeholder else "body"
#            })
#            idx += 1
#        # -----------------------------
#        # GROUP SHAPES
#        # -----------------------------
#        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
#            for shp in shape.shapes:
#                if _is_editable_text_shape(shp):
#                    editable_shapes.append({
#                        "shape_id": f"shape_{idx}",
#                        "text": shp.text.strip(),
#                        "placeholder": False,
#                        "type": "body"
#                    })
#                    idx += 1
#    png_path = export_slide_to_png(ppt_path, slide_index)
#    return {
#        "slide_index": slide_index,
#        "ppt_path": ppt_path,
#        "png_path": png_path,   # None if preview fails
#        "editable_shapes": editable_shapes
#    }

# # =========================================================
# # PREVIEW EXPORT (OS AWARE)
# # =========================================================
# def export_slide_to_png(ppt_path, slide_index):
#    """
#    Export slide to PNG.
#    - Windows → PowerPoint COM
#    - Linux/Docker → LibreOffice headless
#    """
#    try:
#        if IS_WINDOWS:
#            return _export_with_powerpoint(ppt_path, slide_index)
#        return _export_with_libreoffice(ppt_path, slide_index)
#    except Exception as e:
#        # Preview should NEVER crash the app
#        print(f"[WARN] Slide preview failed: {e}")
#        return None

# # =========================================================
# # WINDOWS IMPLEMENTATION (PowerPoint COM)
# # =========================================================
# def _export_with_powerpoint(ppt_path, slide_index):
#    import pythoncom
#    import win32com.client
#    pythoncom.CoInitialize()
#    try:
#        ppt = win32com.client.Dispatch("PowerPoint.Application")
#        ppt.Visible = True
#        pres = ppt.Presentations.Open(ppt_path, WithWindow=True)
#        slide = pres.Slides[slide_index + 1]  # PowerPoint is 1-based
#        out_path = os.path.join(
#            os.path.dirname(ppt_path),
#            f"slide_{slide_index}_{uuid.uuid4().hex[:6]}.png"
#        )
#        slide.Export(out_path, "PNG", 1920, 1080)
#        pres.Close()
#        ppt.Quit()
#        return out_path
#    finally:
#        pythoncom.CoUninitialize()

# # =========================================================
# # LINUX / DOCKER IMPLEMENTATION (LibreOffice)
# # =========================================================
# def _export_with_libreoffice(ppt_path, slide_index):
#    output_dir = os.path.dirname(ppt_path)
#    base_name = os.path.splitext(os.path.basename(ppt_path))[0]
#    libreoffice_cmd = "libreoffice"
#    if not shutil.which(libreoffice_cmd):
#        libreoffice_cmd = "soffice"
#    subprocess.run(
#        [
#            libreoffice_cmd,
#            "--headless",
#            "--convert-to", "png",
#            "--outdir", output_dir,
#            ppt_path
#        ],
#        check=True
#    )
#    # LibreOffice produces ONE PNG only
#    png_path = os.path.join(output_dir, f"{base_name}.png")
#    return png_path if os.path.exists(png_path) else None

# # =========================================================
# # HELPERS
# # =========================================================
# def _is_editable_text_shape(shape):
#    if not shape.has_text_frame:
#        return False
#    text = shape.text.strip()
#    if not text or len(text) < 3:
#        return False
#    return True



import os
import uuid
import platform
import subprocess
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
# =========================================================
# OS DETECTION
# =========================================================
IS_WINDOWS = platform.system() == "Windows"

# =========================================================
# PUBLIC API (USED BY FRONTEND)
# =========================================================
def extract_slide_structure(ppt_path: str, slide_index: int) -> dict:
   """
   Extract editable text shapes + slide preview image.
   - Text extraction: python-pptx (cross-platform)
   - Preview:
       Windows → PowerPoint COM
       Linux/Docker → PPT → PDF → PNG (LibreOffice + poppler)
   Preview failure NEVER crashes the app.
   """
   prs = Presentation(ppt_path)
   slide = prs.slides[slide_index]
   editable_shapes = []
   idx = 0
   for shape in slide.shapes:
       # -----------------------------
       # TEXT SHAPES
       # -----------------------------
       if _is_editable_text_shape(shape):
           is_placeholder = False
           try:
               _ = shape.placeholder_format
               is_placeholder = True
           except Exception:
               pass
           editable_shapes.append({
               "shape_id": f"shape_{idx}",
               "text": shape.text.strip(),
               "placeholder": is_placeholder,
               "type": "title" if is_placeholder else "body"
           })
           idx += 1
       # -----------------------------
       # GROUP SHAPES
       # -----------------------------
       elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
           for shp in shape.shapes:
               if _is_editable_text_shape(shp):
                   editable_shapes.append({
                       "shape_id": f"shape_{idx}",
                       "text": shp.text.strip(),
                       "placeholder": False,
                       "type": "body"
                   })
                   idx += 1
   png_path = export_slide_to_png(ppt_path, slide_index)
   return {
       "slide_index": slide_index,
       "ppt_path": ppt_path,
       "png_path": png_path,  # None if preview fails
       "editable_shapes": editable_shapes
   }

# =========================================================
# PREVIEW EXPORT (OS-AWARE)
# =========================================================
def export_slide_to_png(ppt_path: str, slide_index: int) -> str | None:
   """
   Export ONE slide to PNG.
   Never throws – preview failure should not crash UI.
   """
   try:
       if IS_WINDOWS:
           return _export_with_powerpoint(ppt_path, slide_index)
       return _export_with_libreoffice(ppt_path, slide_index)
   except Exception as e:
       print(f"[WARN] Slide preview failed: {e}")
       return None

# =========================================================
# WINDOWS IMPLEMENTATION (PowerPoint COM)
# =========================================================
def _export_with_powerpoint(ppt_path: str, slide_index: int) -> str:
   import pythoncom
   import win32com.client
   pythoncom.CoInitialize()
   try:
       ppt = win32com.client.Dispatch("PowerPoint.Application")
       ppt.Visible = True
       pres = ppt.Presentations.Open(ppt_path, WithWindow=True)
       slide = pres.Slides[slide_index + 1]  # PowerPoint is 1-based
       out_path = os.path.join(
           os.path.dirname(ppt_path),
           f"slide_{slide_index}_{uuid.uuid4().hex[:6]}.png"
       )
       slide.Export(out_path, "PNG", 1920, 1080)
       pres.Close()
       ppt.Quit()
       return out_path
   finally:
       pythoncom.CoUninitialize()

# =========================================================
# LINUX / DOCKER IMPLEMENTATION (PPT → PDF → PNG)
# =========================================================
def _export_with_libreoffice(ppt_path: str, slide_index: int) -> str | None:
   """
   Reliable Linux pipeline:
       PPT → PDF (LibreOffice)
       PDF → PNG (pdf2image)
   Guarantees correct slide ordering.
   """
   output_dir = os.path.dirname(ppt_path)
   base_name = os.path.splitext(os.path.basename(ppt_path))[0]
   pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
   # ---------- PPT → PDF ----------
   libreoffice_cmd = "libreoffice"
   if not shutil.which(libreoffice_cmd):
       libreoffice_cmd = "soffice"
   if not os.path.exists(pdf_path):
       subprocess.run(
           [
               libreoffice_cmd,
               "--headless",
               "--convert-to", "pdf",
               "--outdir", output_dir,
               ppt_path
           ],
           check=True
       )
   # ---------- PDF → PNG ----------
   from pdf2image import convert_from_path
   images = convert_from_path(pdf_path, dpi=200)
   if slide_index >= len(images):
       return None
   out_path = os.path.join(
       output_dir,
       f"{base_name}_slide_{slide_index}.png"
   )
   images[slide_index].save(out_path, "PNG")
   return out_path

# =========================================================
# HELPERS
# =========================================================
def _is_editable_text_shape(shape) -> bool:
   if not shape.has_text_frame:
       return False
   text = shape.text.strip()
   if not text or len(text) < 3:
       return False
   return True