
import os

import uuid

import platform

import subprocess

import shutil

from pptx import Presentation

from pptx.enum.shapes import MSO_SHAPE_TYPE

# =========================================================

# OS DETECTION

# =========================================================

IS_WINDOWS = platform.system() == "Windows"

# =========================================================

# WINDOWS POWERPOINT COM GLOBALS (IMPORTANT)

# PowerPoint MUST be opened once per PPT, not per slide

# =========================================================

_ppt_app = None

_ppt_pres = None

_ppt_path = None

# =========================================================

# PUBLIC API

# =========================================================

def extract_slide_structure(ppt_path: str, slide_index: int) -> dict:

    """

    Extract editable text shapes + slide preview image.

    Windows:

        PowerPoint COM (opened once, reused across slides)

    Linux / Docker:

        LibreOffice -> PDF -> PNG (stateless, unchanged)

    """

    prs = Presentation(ppt_path)

    slide = prs.slides[slide_index]

    editable_shapes = []

    idx = 0

    for shape in slide.shapes:

        # -----------------------------

        # TEXT SHAPES

        # -----------------------------

        if _is_editable_text_shape(shape):

            is_placeholder = False

            try:

                _ = shape.placeholder_format

                is_placeholder = True

            except Exception:

                pass

            editable_shapes.append({

                "shape_id": f"shape_{idx}",

                "text": shape.text.strip(),

                "placeholder": is_placeholder,

                "type": "title" if is_placeholder else "body"

            })

            idx += 1

        # -----------------------------

        # GROUP SHAPES

        # -----------------------------

        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:

            for shp in shape.shapes:

                if _is_editable_text_shape(shp):

                    editable_shapes.append({

                        "shape_id": f"shape_{idx}",

                        "text": shp.text.strip(),

                        "placeholder": False,

                        "type": "body"

                    })

                    idx += 1

    png_path = export_slide_to_png(ppt_path, slide_index)

    return {

        "slide_index": slide_index,

        "ppt_path": ppt_path,

        "png_path": png_path,

        "editable_shapes": editable_shapes

    }

# =========================================================

# PREVIEW EXPORT DISPATCHER

# =========================================================

def export_slide_to_png(ppt_path: str, slide_index: int) -> str | None:

    if IS_WINDOWS:

        return _export_with_powerpoint_windows(ppt_path, slide_index)

    return _export_with_libreoffice(ppt_path, slide_index)

# =========================================================

# WINDOWS IMPLEMENTATION (FIXED & LOOP-SAFE)

# =========================================================

def _export_with_powerpoint_windows(ppt_path: str, slide_index: int) -> str:

    """

    IMPORTANT:

    - PowerPoint is opened ONCE per PPT

    - Presentation is NEVER closed inside the slide loop

    - This avoids Presentation.Close() COM failures

    """

    global _ppt_app, _ppt_pres, _ppt_path

    import pythoncom

    import win32com.client

    pythoncom.CoInitialize()

    # Open PowerPoint once per PPT

    if _ppt_app is None or _ppt_path != ppt_path:

        _ppt_app = win32com.client.Dispatch("PowerPoint.Application")

        _ppt_app.Visible = True

        _ppt_pres = _ppt_app.Presentations.Open(ppt_path, WithWindow=True)

        _ppt_path = ppt_path

    slide = _ppt_pres.Slides[slide_index]

    out_path = os.path.join(

        os.path.dirname(ppt_path),

        f"slide_{slide_index}_{uuid.uuid4().hex[:6]}.png"

    )

    slide.Export(out_path, "PNG", 1920, 1080)

    return out_path

# =========================================================

# OPTIONAL CLEANUP (CALL ON APP SHUTDOWN IF NEEDED)

# =========================================================

def cleanup_powerpoint():

    global _ppt_app, _ppt_pres, _ppt_path

    try:

        if _ppt_pres:

            _ppt_pres.Close()

        if _ppt_app:

            _ppt_app.Quit()

    except Exception:

        pass

    _ppt_app = None

    _ppt_pres = None

    _ppt_path = None

# =========================================================

# LINUX / DOCKER IMPLEMENTATION (UNCHANGED)

# =========================================================

def _export_with_libreoffice(ppt_path: str, slide_index: int) -> str | None:

    """

    Reliable Linux pipeline:

        PPT → PDF (LibreOffice)

        PDF → PNG (pdf2image)

    """

    output_dir = os.path.dirname(ppt_path)

    base_name = os.path.splitext(os.path.basename(ppt_path))[0]

    pdf_path = os.path.join(output_dir, f"{base_name}.pdf")

    libreoffice_cmd = "libreoffice"

    if not shutil.which(libreoffice_cmd):

        libreoffice_cmd = "soffice"

    # Always regenerate to avoid stale previews

    if os.path.exists(pdf_path):

        os.remove(pdf_path)

    subprocess.run(

        [

            libreoffice_cmd,

            "--headless",

            "--convert-to", "pdf",

            "--outdir", output_dir,

            ppt_path

        ],

        check=True

    )

    from pdf2image import convert_from_path

    images = convert_from_path(pdf_path, dpi=200)

    if slide_index >= len(images):

        return None

    out_path = os.path.join(

        output_dir,

        f"{base_name}_slide_{slide_index}.png"

    )

    images[slide_index].save(out_path, "PNG")

    return out_path

# =========================================================

# HELPERS

# =========================================================

def _is_editable_text_shape(shape) -> bool:

    if not shape.has_text_frame:

        return False

    text = shape.text.strip()

    if not text or len(text) < 3:

        return False

    return True
 
 
 
