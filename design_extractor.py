import os
import json
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from azure.storage.blob import BlobServiceClient
from utils import get_env, logger, now_ts

# === CONFIG ===
BLOB_CONN = get_env("AZURE_BLOB_CONN", required=True)
BLOB_CONTAINER = get_env("AZURE_BLOB_CONTAINER", "ppt-dataset")

# === BLOB CLIENT ===
blob_service = BlobServiceClient.from_connection_string(BLOB_CONN)
container_client = blob_service.get_container_client(BLOB_CONTAINER)

# === LOCAL OUTPUT DIR ===
OUTPUT_DIR = "design_jsons"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def ppt_already_processed(ppt_name: str) -> bool:
    """Check if design JSON already exists locally."""
    json_path = os.path.join(OUTPUT_DIR, f"{os.path.basename(ppt_name)}.json")
    return os.path.exists(json_path)


def extract_theme_colors(prs):
    """Extract the theme color palette from the presentation."""
    colors = []
    try:
        if prs.theme and prs.theme.color_scheme:
            for attr in dir(prs.theme.color_scheme):
                if not attr.startswith("_"):
                    color = getattr(prs.theme.color_scheme, attr)
                    if hasattr(color, "rgb") and color.rgb:
                        colors.append(f"#{color.rgb}")
    except Exception:
        pass
    return list(set(colors))


def extract_fonts_and_layouts(prs):
    """Extract fonts and layout names used across slides."""
    fonts = set()
    layouts = set()
    try:
        for slide in prs.slides:
            if slide.slide_layout and slide.slide_layout.name:
                layouts.add(slide.slide_layout.name)
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font and r.font.name:
                                fonts.add(r.font.name)
    except Exception:
        pass
    return list(fonts), list(layouts)


def extract_slide_elements(slide):
    """Extract background color, shapes, charts, tables, and icons for each slide."""
    slide_data = {
        "layout": None,
        "background_color": None,
        "images": [],
        "shapes": [],
        "charts": [],
        "tables": [],
        "icons": [],
        "text_fonts": [],
        "accent_colors": [],
    }

    # Layout name
    try:
        if slide.slide_layout and slide.slide_layout.name:
            slide_data["layout"] = slide.slide_layout.name
    except Exception:
        pass

    # Background color
    try:
        bg = slide.background.fill
        if bg and bg.fore_color and hasattr(bg.fore_color, "rgb"):
            slide_data["background_color"] = f"#{bg.fore_color.rgb}"
    except Exception:
        pass

    # Shapes and elements
    for shape in slide.shapes:
        stype = shape.shape_type
        try:
            if stype == MSO_SHAPE_TYPE.PICTURE:
                slide_data["images"].append(shape.name)
            elif stype == MSO_SHAPE_TYPE.TABLE:
                slide_data["tables"].append(shape.name)
            elif stype == MSO_SHAPE_TYPE.CHART:
                try:
                    slide_data["charts"].append(str(shape.chart.chart_type))
                except Exception:
                    slide_data["charts"].append("Chart")
            elif stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
                slide_data["shapes"].append(str(shape.auto_shape_type))
            elif "icon" in shape.name.lower():
                slide_data["icons"].append(shape.name)

            # Accent colors from shape fill or outline
            if hasattr(shape, "fill") and shape.fill and shape.fill.fore_color and hasattr(shape.fill.fore_color, "rgb"):
                slide_data["accent_colors"].append(f"#{shape.fill.fore_color.rgb}")
            if hasattr(shape, "line") and shape.line and shape.line.color and hasattr(shape.line.color, "rgb"):
                slide_data["accent_colors"].append(f"#{shape.line.color.rgb}")

            # Text fonts
            if hasattr(shape, "text_frame") and shape.text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font and r.font.name:
                            slide_data["text_fonts"].append(r.font.name)

        except Exception:
            logger.debug(f"Skipped shape on slide due to parsing issue.")
            continue

    # Deduplicate lists
    for key in ["accent_colors", "text_fonts"]:
        slide_data[key] = list(set(slide_data[key]))

    return slide_data


def extract_design_elements(local_path: str):
    """Extract design metadata for a PPT."""
    prs = Presentation(local_path)
    theme_colors = extract_theme_colors(prs)
    fonts, layouts = extract_fonts_and_layouts(prs)

    design_data = {
        "ppt_name": os.path.basename(local_path),
        "extracted_on": now_ts(),
        "theme_colors": theme_colors,
        "font_families": fonts,
        "slide_layouts": layouts,
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_data = extract_slide_elements(slide)
        slide_data["index"] = i
        slide_data["slide_number"] = i + 1 # ✅ Added slide number (1-based)
        design_data["slides"].append(slide_data)

    return design_data


def process_blob(blob_name: str):
    """Download PPT from Azure Blob, extract design info, and save as JSON."""
    json_path = os.path.join(OUTPUT_DIR, f"{os.path.basename(blob_name)}.json")

    if ppt_already_processed(blob_name):
        logger.info(f"Skipping {blob_name} — design JSON already exists.")
        return

    logger.info(f"Processing design extraction for: {blob_name}")
    tmp_path = os.path.join(tempfile.gettempdir(), blob_name.replace("/", "_"))

    try:
        with open(tmp_path, "wb") as fp:
            stream = container_client.download_blob(blob_name)
            stream.readinto(fp)
    except Exception as e:
        logger.error(f"Failed to download {blob_name} from blob: {e}")
        return

    try:
        design_data = extract_design_elements(tmp_path)
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(design_data, f, indent=2)
        logger.info(f"✅ Saved design JSON: {json_path}")
    except Exception as e:
        logger.exception(f"Design extraction failed for {blob_name}: {e}")


def main():
    """Main function: iterate over blob PPT files and extract design JSONs."""
    logger.info("Starting design extraction from Azure Blob...")
    try:
        for blob in container_client.list_blobs():
            if blob.name.endswith(".pptx") or blob.name.endswith(".ppt"):
                process_blob(blob.name)
    except Exception as e:
        logger.exception(f"Error listing blobs or processing files: {e}")
    logger.info("Design extraction complete.")


if __name__ == "__main__":
    main()
