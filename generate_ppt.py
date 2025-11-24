# ============================================================
# generate_ppt.py – FINAL FIXED WITH NORMALIZED PPT NAME MATCH
# ============================================================

import os
import tempfile
import uuid
import json
import re
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from utils import (
    get_env, safe_json_load, logger, now_ts,
    ensure_dir, text_client, image_client
)
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


# ------------------------------------------------------------
# NORMALIZE PPT NAME (CRITICAL FIX)
# ------------------------------------------------------------
def normalize_name(path):
    """
    Normalize any PPT path from Chroma or Blob:
    - Replace backslashes
    - Strip OneDrive / full paths
    - Extract only the final filename
    """
    if not path:
        return ""
    return path.replace("\\", "/").split("/")[-1]


# ------------------------------------------------------------
# INITIAL SETUP
# ------------------------------------------------------------

ensure_dir("design_jsons")

CHAT_MODEL = get_env("CHAT_MODEL", required=True)
IMAGE_MODEL = get_env("IMAGE_MODEL", required=True)


# ------------------------------------------------------------
# USER INTENT PARSER
# ------------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    match = re.search(r"(\d+)\s+slides?", prompt.lower())
    if match:
        num_slides = int(match.group(1))

    for t in ["corporate", "modern", "minimal", "professional", "dark", "light"]:
        if t in prompt.lower():
            theme = t.capitalize()
            break

    return num_slides, theme


# ------------------------------------------------------------
# GPT PLAN GENERATOR
# ------------------------------------------------------------
def call_llm_plan(prompt, style, design_context, references_text,
                  num_slides=None, theme=None):

    sys_prompt = (
        "You are a senior presentation designer.\n"
        "Return STRICT JSON ONLY:\n"
        "[{\"title\": str, \"bullets\": [str], \"visual_required\": bool, \"visual_prompt\": str }]\n"
        "If user mentions images → set visual_required=true.\n"
        "Do not include text inside image prompts.\n"
        "Use design cues only when appropriate.\n"
    )

    user_prompt = f"Create a {style} presentation plan for: {prompt}"
    if num_slides:
        user_prompt += f". Make exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ],
            max_completion_tokens=1500,
            temperature=1
        )
        plan = safe_json_load(resp.choices[0].message.content)

        if not plan:
            logger.warning("Invalid JSON, using fallback slide.")
            return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]

        return plan

    except Exception:
        logger.exception("Plan generation failed")
        return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]


# ------------------------------------------------------------
# IMAGE GENERATION
# ------------------------------------------------------------
def generate_visual_image(prompt: str):
    img_prompt = (prompt or "") + " Minimal professional illustration. No text labels."

    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=img_prompt,
            size="1024x1024"
        )

        b64 = getattr(resp.data[0], "b64_json", None)
        if b64:
            img_bytes = base64.b64decode(b64)
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tmp.write(img_bytes)
            tmp.close()
            return tmp.name

        url = getattr(resp.data[0], "url", None)
        if url:
            import requests
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tmp.write(requests.get(url).content)
            tmp.close()
            return tmp.name

        return None

    except Exception:
        logger.exception("Image generation failed")
        return None


# ------------------------------------------------------------
# STYLE HELPERS
# ------------------------------------------------------------
def parse_hex_color(hex_str):
    try:
        if not hex_str:
            return None
        hex_str = hex_str.replace("#", "")
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return RGBColor(r, g, b)
    except:
        return None


# ------------------------------------------------------------
# EXTRACT DESIGN FOR SLIDE (WITH NAME NORMALIZATION FIX)
# ------------------------------------------------------------
def extract_design_for_slide(design_context, matched_ppt_name, matched_slide_idx):

    norm_target = normalize_name(matched_ppt_name)
    logger.info(f"[EXTRACT] Looking for PPT={norm_target}, slide={matched_slide_idx}")

    for design in design_context:
        norm_json = normalize_name(design["ppt_name"])

        logger.info(f"[EXTRACT] Checking JSON PPT={norm_json}")

        if norm_json == norm_target:
            logger.info("[EXTRACT] PPT name matched!")

            for slide in design["slides"]:
                if slide["index"] == matched_slide_idx:
                    logger.info("[EXTRACT] Slide index matched. RETURNING.")
                    return slide

            logger.info("[EXTRACT] Slide index NOT matched")
        else:
            logger.info("[EXTRACT] PPT name NOT matched")

    logger.info("[EXTRACT] No matching PPT found in design context")
    return None


# ------------------------------------------------------------
# APPLY DESIGN STYLE (Option 2 + Option 3)
# ------------------------------------------------------------
def apply_design_style(slide, prs, design_meta):

    if not design_meta:
        return

    # Font family
    text_fonts = design_meta.get("text_fonts", [])
    if text_fonts:
        chosen = text_fonts[0]
        try:
            body = slide.placeholders[1]
            for p in body.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = chosen
        except:
            pass

    # Accent color
    accents = design_meta.get("accent_colors", [])
    if accents:
        accent_color = parse_hex_color(accents[0])

    # Background color
    bg = parse_hex_color(design_meta.get("background_color"))
    if bg:
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg
        except:
            pass


# ------------------------------------------------------------
# PPT BUILDER
# ------------------------------------------------------------
def build_ppt(slides, matched_designs):
    prs = Presentation()

    for idx, s in enumerate(slides):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Apply design style
        apply_design_style(slide, prs, matched_designs[idx])

        # Title
        slide.shapes.title.text = s.get("title", "")

        # Body text
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for b in s.get("bullets", []):
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)

        # Move text slightly down
        body.top += Inches(0.25)

        # No image → full width text
        if not s.get("image_path"):
            continue

        # Shrink text to left side
        body.width = prs.slide_width - Inches(4.0)

        # Insert right-side image
        try:
            img_path = s["image_path"]
            img = Image.open(img_path)
            w, h = img.size
            aspect = w / h

            max_w = Inches(3.0)
            max_h = Inches(2.5)

            if aspect >= 1:
                final_w = max_w
                final_h = final_w / aspect
            else:
                final_h = max_h
                final_w = final_h * aspect

            left = prs.slide_width - final_w - Inches(0.5)
            top = slide.shapes.title.top + slide.shapes.title.height + Inches(0.2)

            slide.shapes.add_picture(img_path, left, top, width=final_w, height=final_h)

        except Exception:
            logger.exception("Image placement failed")

    out_path = os.path.join(tempfile.gettempdir(), f"generated_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(out_path)
    return out_path


# ------------------------------------------------------------
# MAIN PIPELINE
# ------------------------------------------------------------
def generate_presentation(prompt: str, style="Auto", requested_num_slides=None,
                          theme=None, tag_filters=None):

    detected_slides, detected_theme = parse_user_intent(prompt)
    requested_num_slides = requested_num_slides or detected_slides
    theme = theme or detected_theme

    refs = semantic_search(prompt, top_k=5, tags=tag_filters) or []

    design_context = []
    reference_text = []

    # Load design JSONs
    for r in refs:
        ppt_name = r.get("ppt_name")
        snippet = (r.get("text") or "")[:500]
        if snippet:
            reference_text.append(snippet)

        json_path = os.path.join("design_jsons", normalize_name(ppt_name) + ".json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_text)} text snippets.")

    # LLM plan
    plan = call_llm_plan(
        prompt, style, design_context, reference_text,
        num_slides=requested_num_slides, theme=theme
    )

    force_images = "image" in prompt.lower() or "images" in prompt.lower()

    slides = []
    matched_designs = []

    # PER-SLIDE DESIGN MATCHING
    for sp in plan:

        best = refs[0] if refs else None
        design_meta = None

        if best:
            logger.info(f"[MATCH] Trying PPT={best.get('ppt_name')}, index={best.get('slide_index')}")

            try:
                matched_idx = int(best.get("slide_index"))
            except:
                matched_idx = 0

            design_meta = extract_design_for_slide(
                design_context,
                best.get("ppt_name"),
                matched_idx
            )

            logger.info(f"[MATCH] Design match result: {design_meta is not None}")

        matched_designs.append(design_meta)

        img = None
        if sp.get("visual_required") or force_images:
            img = generate_visual_image(sp.get("visual_prompt"))

        slides.append({
            "title": sp.get("title"),
            "bullets": sp.get("bullets", []),
            "image_path": img
        })

    out_path = build_ppt(slides, matched_designs)

    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(out_path, fname)

    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "slides_generated": len(slides),
        "ppt_file": fname
    }

    upload_json_to_blob(json.dumps(log, indent=2).encode(), f"logs/{fname}.json")

    return out_path, log